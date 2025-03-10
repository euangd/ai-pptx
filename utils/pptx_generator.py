import json
import random
import re
import time
import logging
import pptx
from pptx import Presentation

from utils.llm import LLM
from utils.ppt_tools import recreate_slide_by_win32
from utils.prompter import PromptLibrarian

logger = logging.getLogger(__name__)


class PptxGenerator:
    PPT_PARAM_PATTERN = r'\{(.*?)\}'
    MD_CODE_JSON_PATTERN = r'```json.*?\n(.*?)```'

    def __init__(self, llm: LLM, save_path: str, template_path: str = None):
        self.llm = llm
        self.save_path = save_path
        self.has_template = False
        if template_path:
            self.has_template = True
            self.template_path = template_path
            self.template_ppt = Presentation(template_path)
            self.template_params = self._extract_params_from_template()

    def _extract_params_from_template(self):
        """Generate PPT file"""
        # Note that when extracting parameters, you need to unlock all <groups> in the template PPT, otherwise you may not find text boxes.
        start_slide_idx = 0
        catalogue_slide_idx = 1
        title_slide_idx = 2
        content_slide_idxs = [3, 4, 5, 6, 7, 8]
        end_slide_idx = 9
        template_params = {
            "first_slide": {"nos": [start_slide_idx], "params": []},
            "catalogue_slide": {"nos": [catalogue_slide_idx], "params": []},
            "title_slide": {"nos": [title_slide_idx], "params": []},
            "content_slide": {"nos": content_slide_idxs, "params": []},
            "end_slide": {"nos": [end_slide_idx], "params": []}
        }

        # The {params} on the same page in PPT must be different to avoid confusion, no requirements for different pages
        for slide_name, slide_info in template_params.items():
            nos = slide_info["nos"]
            for n in nos:
                slide = self.template_ppt.slides[n]
                temp_params = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                matches = [match.group(1) for match in re.finditer(self.PPT_PARAM_PATTERN, run.text)]
                                temp_params.extend(matches)
                slide_info["params"].append(temp_params)
        return template_params

    def llm_generate_online_content(self, topic: str):
        """Generate outline based on topic"""
        output_format = json.dumps(
            {
                "topic": "str",
                "pages": [
                    {
                        "title": "str",
                        "pages": [
                            {"sub_title": "str", "desc": "str", "content": "str"}
                        ]
                    }
                ]
            }, ensure_ascii=False
        )
        # remove space to save token
        output_format = output_format.replace(" ", "")
        prompt = PromptLibrarian.read(path="ppt.generate_content.v1").format(topic=topic,
                                                                             language="British English",
                                                                             output_format=output_format)
        messages = [
            {"role": "system", "content": "You are an all-capable assistant"},
            {"role": "user", "content": prompt}
        ]
        c = self.llm.chat_in_all(messages)
        if c[-1] != "}":
            logger.warning("[Continuous] Reply not end, go on ...")
            messages.append({"role": "assistant", "content": c})
            messages.append({"role": "user", "content": "Continue"})
            c += self.llm.chat_in_all(messages)
        return c

    def _llm_generate_content_slide_in_single(self, prompt: str, temperature: float, tp: dict):
        is_match = True
        try_count = 0
        while try_count <= 3:
            try_count += 1
            ctx = self.llm.chat_once(prompt=prompt, temperature=temperature)
            m = re.findall(self.MD_CODE_JSON_PATTERN, ctx, re.DOTALL)
            if m: ctx = m[0].replace("'", '"')
            # try to load json
            try:
                ctx = json.loads(ctx)
            except Exception as e:
                logger.warning(f"ctx json.loads error: \n{ctx}")
                time.sleep(0.8 * try_count)
                continue
            # try to match params
            gcs = [gc for gc in ctx.keys()]
            for tk in tp.keys():
                if tk not in gcs:
                    is_match = False
                    break
            if is_match:
                logger.info(f"try generated count <{try_count}>, ctx: \n{ctx}")
                return ctx
            time.sleep(0.8 * try_count)
        return None

    def llm_generate_content_slide_content(self, topic: str, online_content: str):
        """Generate complete content based on outline"""
        logger.info(f"online_content: \n{online_content}")
        online_content = json.loads(online_content)
        current_online_content = online_content["pages"]
        content_slide = self.template_params["content_slide"]

        # Add title number, subtitle number
        for idx, c in enumerate(current_online_content):
            c["no"] = idx + 1
        for c in current_online_content:
            for idx, s in enumerate(c["pages"]):
                s["sub_no"] = idx + 1

        title_count = len(current_online_content)
        resorted_no_idxs = random.sample(range(len(content_slide["nos"])), k=title_count)
        current_template_resort_nos = [content_slide["nos"][idx] for idx in resorted_no_idxs]
        logger.info(f"content_slide_params: {content_slide['params']}")
        logger.info(f"resorted_no_idxs: {resorted_no_idxs}")
        logger.info(f"current_template_resort_nos: {current_template_resort_nos}")

        current_template_params = [
            {k: "" for k in content_slide["params"][idx]}
            for idx in resorted_no_idxs
        ]

        for oc, tp in zip(current_online_content, current_template_params):
            title = oc["title"]
            prompt = f"""# Info
## OnlineJson
```{oc}```
## TemplateParamsJson
```{tp}```
# Tasks
Strictly follow [Info.TemplateParamsJson], based on the content of the `{title}` title in 《{topic}》, fill in [Info.OnlineJson] accordingly, and finally output according to the markdown json format.
Note: The key values of json strictly correspond to [Info.TemplateParamsJson], and the values corresponding to keys cannot contain lists or dictionaries.
------
output:"""
            ctx = self._llm_generate_content_slide_in_single(prompt, 0.6, tp)
            if ctx:
                # Strictly match and assign values according to template parameters
                for tk in tp.keys():
                    tp[tk] = ctx.get(tk, "")
                time.sleep(2)
            else:
                logger.exception(f"failed to generate content for title: {title}. Skip it!")

        data = {
            "titles_param": {f'title_{i + 1}': c["title"] for i, c in enumerate(current_online_content)},
            "contents_param": current_template_params,
            "nos": current_template_resort_nos
        }
        return data

    def generate_ppt(self, meta_info: dict, generation_content: dict):
        """
        generate ppt based on content which is generated by llm
        :param meta_info:
        :param generation_content:
        :return:
        """
        # 1. Create a new ppt based on template
        logger.info(f"meta_info: {meta_info}")
        logger.info(f"generation_content: {generation_content}")

        titles_param = generation_content["titles_param"]
        contents_param = generation_content["contents_param"]
        nos = generation_content["nos"]

        all_params = contents_param
        # Insert home page content and slide No
        all_params.insert(0, meta_info)
        nos.insert(0, self.template_params["first_slide"]["nos"][0])
        # Insert directory page
        all_params.insert(1, titles_param)
        nos.insert(1, self.template_params["catalogue_slide"]["nos"][0])
        # Insert end page (at the end)
        all_params.append({})
        nos.append(self.template_params["end_slide"]["nos"][0])

        # 2. Regenerate ppt based on rearranged slide No
        logger.info(f"nos: {nos}")
        recreate_slide_by_win32(self.template_path, self.save_path, indexs=nos)

        # 3. Add pages and fill in content on new ppt
        new_ppt = pptx.Presentation(self.save_path)
        for idx, p_dict in enumerate(all_params):
            for shape in new_ppt.slides[idx].shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            for match in re.findall(pattern=self.PPT_PARAM_PATTERN, string=run.text):
                                m_str = "{" + match + "}"
                                m_key = match
                                run.text = run.text.replace(m_str, str(p_dict.get(m_key, '')))

        # 4. Save PPT
        new_ppt.save(self.save_path)

    def generate(self, meta_info: dict):
        """Generate PPT according to template"""
        online_content = self.llm_generate_online_content(meta_info["topic"])
        generation_content = self.llm_generate_content_slide_content(meta_info["topic"], online_content)
        self.generate_ppt(meta_info, generation_content)
